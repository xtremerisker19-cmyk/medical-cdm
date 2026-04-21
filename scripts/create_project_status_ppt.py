from __future__ import annotations

from collections import Counter
from datetime import date
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


WORKBOOK_PATH = Path("classification 1.xlsx")
OUTPUT_PATH = Path("Project_Status_Update.pptx")


def normalize(value: object) -> str:
    if value is None:
        return ""
    return str(value).strip()


def parse_sheet(worksheet) -> tuple[list[str], list[dict[str, str]]]:
    headers = [normalize(cell.value) for cell in worksheet[1]]
    rows: list[dict[str, str]] = []
    for source_row in worksheet.iter_rows(min_row=2, values_only=True):
        if all(normalize(cell) == "" for cell in source_row):
            continue
        row_dict: dict[str, str] = {}
        for index, header in enumerate(headers):
            key = header if header else f"COL_{index + 1}"
            value = source_row[index] if index < len(source_row) else ""
            row_dict[key] = normalize(value)
        rows.append(row_dict)
    return headers, rows


def get_column_key(headers: list[str], starts_with: str) -> str | None:
    normalized_prefix = starts_with.lower()
    for header in headers:
        if header.lower().startswith(normalized_prefix):
            return header
    return None


def collect_metrics(workbook_path: Path) -> dict[str, object]:
    workbook = openpyxl.load_workbook(workbook_path, data_only=True)

    parsed: dict[str, dict[str, object]] = {}
    for sheet in workbook.worksheets:
        headers, rows = parse_sheet(sheet)
        parsed[sheet.title] = {"headers": headers, "rows": rows}

    source_counts = {
        sheet_name: len(data["rows"]) for sheet_name, data in parsed.items()  # type: ignore[index]
    }
    total_records = sum(source_counts.values())

    dashboard_counts: Counter[str] = Counter()
    category_counts: Counter[str] = Counter()

    unknown_markers = {"?", "not found", "nothing found", "nothing in dataflow"}
    unknown_cell_count = 0
    blank_cell_count = 0

    for data in parsed.values():
        headers = data["headers"]  # type: ignore[index]
        rows = data["rows"]  # type: ignore[index]

        category_key = get_column_key(headers, "category")
        dashboard_key = None
        for header in headers:
            if "dash" in header.lower():
                dashboard_key = header
                break

        for row in rows:
            if category_key:
                category = row.get(category_key, "")
                if category:
                    category_counts[category] += 1
            if dashboard_key:
                dashboard = row.get(dashboard_key, "")
                if dashboard:
                    dashboard_counts[dashboard] += 1

            for value in row.values():
                normalized_value = value.lower()
                if value == "":
                    blank_cell_count += 1
                if (
                    normalized_value in unknown_markers
                    or "nothing found" in normalized_value
                    or normalized_value == "?"
                ):
                    unknown_cell_count += 1

    tables_data = parsed.get("TABLES", {})
    tables_headers = tables_data.get("headers", [])
    tables_rows = tables_data.get("rows", [])
    actor_key = (
        get_column_key(tables_headers, "actor classification")
        if isinstance(tables_headers, list)
        else None
    )
    actor_counts: Counter[str] = Counter()
    if actor_key and isinstance(tables_rows, list):
        for row in tables_rows:
            actor = row.get(actor_key, "")
            if actor:
                actor_counts[actor] += 1

    return {
        "sheet_names": list(parsed.keys()),
        "source_counts": source_counts,
        "total_records": total_records,
        "dashboard_counts": dashboard_counts,
        "category_counts": category_counts,
        "actor_counts": actor_counts,
        "unknown_cell_count": unknown_cell_count,
        "blank_cell_count": blank_cell_count,
    }


def style_title(shape, size: int = 34) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(31, 56, 100)


def style_subtitle(shape, size: int = 18) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = RGBColor(68, 84, 106)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Project Status Update"
    style_title(slide.shapes.title, 40)

    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Data inventory snapshot generated from classification workbook\n"
        f"As of {date.today().isoformat()}"
    )
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(68, 84, 106)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 32)

    body = slide.shapes.placeholders[1]
    body.text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.text_frame.paragraphs[0] if index == 0 else body.text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)


def add_bar_chart_slide(
    prs: Presentation,
    title: str,
    categories: list[str],
    values: list[int],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 30)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Count", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(11.6),
        Inches(5.1),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(11)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 30)

    row_count = len(rows) + 1
    col_count = len(headers)
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(0.5),
        Inches(1.5),
        Inches(12.3),
        Inches(5.2),
    )
    table = table_shape.table

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 56, 100)

    for row_index, row_values in enumerate(rows, start=1):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = RGBColor(45, 45, 45)


def build_presentation(metrics: dict[str, object], output_path: Path) -> None:
    prs = Presentation()

    source_counts: dict[str, int] = metrics["source_counts"]  # type: ignore[assignment]
    dashboard_counts: Counter[str] = metrics["dashboard_counts"]  # type: ignore[assignment]
    category_counts: Counter[str] = metrics["category_counts"]  # type: ignore[assignment]
    actor_counts: Counter[str] = metrics["actor_counts"]  # type: ignore[assignment]

    add_title_slide(prs)

    top_source = max(source_counts.items(), key=lambda item: item[1])
    top_dashboard = dashboard_counts.most_common(1)[0] if dashboard_counts else ("N/A", 0)
    top_category = category_counts.most_common(1)[0] if category_counts else ("N/A", 0)
    known_dashboards = len(dashboard_counts)
    known_categories = len(category_counts)

    add_bullet_slide(
        prs,
        "Executive Summary",
        [
            f"Total inventoried objects: {metrics['total_records']} across {len(source_counts)} source layers.",
            f"Largest source layer: {top_source[0]} ({top_source[1]} objects).",
            f"Dashboard footprint: {known_dashboards} dashboards, top is {top_dashboard[0]} ({top_dashboard[1]} objects).",
            f"Category coverage: {known_categories} categories, top is {top_category[0]} ({top_category[1]} objects).",
            f"Data quality watchpoints: {metrics['unknown_cell_count']} unknown markers and {metrics['blank_cell_count']} blank cells.",
        ],
    )

    source_labels = list(source_counts.keys())
    source_values = list(source_counts.values())
    add_bar_chart_slide(prs, "Inventory by Source Layer", source_labels, source_values)

    top_dashboards = dashboard_counts.most_common(8)
    add_bar_chart_slide(
        prs,
        "Top Dashboard Coverage",
        [item[0] for item in top_dashboards],
        [item[1] for item in top_dashboards],
    )

    top_categories = category_counts.most_common(8)
    add_bar_chart_slide(
        prs,
        "Top Category Distribution",
        [item[0] for item in top_categories],
        [item[1] for item in top_categories],
    )

    top_actors = actor_counts.most_common(8)
    add_table_slide(
        prs,
        "Actor Classification Snapshot (TABLES sheet)",
        ["Actor Classification", "Count"],
        [[name, str(count)] for name, count in top_actors] or [["No actor data", "0"]],
    )

    add_bullet_slide(
        prs,
        "Recommended Next Actions",
        [
            "Standardize unknown labels (for example '?', 'nothing found') with controlled vocabulary.",
            "Fill blank metadata fields for table description and ownership before governance sign-off.",
            "Prioritize remediation for high-usage dashboards: TRADE, SECONDARY SALES, JOINT CALL.",
            "Define monthly refresh checks to keep source and dashboard mappings current.",
        ],
    )

    prs.save(output_path)


def main() -> None:
    if not WORKBOOK_PATH.exists():
        raise FileNotFoundError(f"Workbook not found: {WORKBOOK_PATH}")

    metrics = collect_metrics(WORKBOOK_PATH)
    build_presentation(metrics, OUTPUT_PATH)
    print(f"Generated presentation: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
